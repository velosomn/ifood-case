"""Gera o deck de 5 slides (presentation/ifood_case.pptx).

Narrativa: problema (cupom pago sem efeito) -> quanto vale um envio ->
o efeito varia e dá para prever -> impacto em dinheiro -> a grande conclusão.
"""
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "presentation" / "figures"
OUT = ROOT / "presentation" / "ifood_case.pptx"

RED = RGBColor(0xEA, 0x1D, 0x2C)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

# ------------------------------------------------ figura do funil (com ordem)
plt.rcParams["figure.dpi"] = 110
seg = [("Viu e usou o cupom", 38.8, "#2E7D32"),
       ("Viu e não usou", 31.2, "#F9A825"),
       ("Nem viu, nem usou", 13.7, "#9E9E9E"),
       ("Usou SEM ter visto", 9.4, "#EA1D2C"),
       ("Usou e viu só depois", 6.9, "#C62828")]
fig, ax = plt.subplots(figsize=(7.2, 3.2))
labels = [s[0] for s in seg][::-1]
vals = [s[1] for s in seg][::-1]
cols = [s[2] for s in seg][::-1]
ax.barh(labels, vals, color=cols)
for i, v in enumerate(vals):
    ax.text(v + 0.5, i, f"{v:.0f}%", va="center", fontsize=10)
ax.set_xlabel("% das ofertas enviadas (BOGO e desconto)")
ax.set_title("O que acontece com cada oferta enviada")
ax.set_xlim(0, 46)
plt.tight_layout()
plt.savefig(FIG / "10_funnel.png", bbox_inches="tight")
plt.close(fig)

# ------------------------------------------------------------------ helpers
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, left, top, width, height, text, size, color=DARK, bold=False,
             align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def band(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()


def content_slide(kicker, title):
    s = prs.slides.add_slide(BLANK)
    band(s)
    add_text(s, 0.7, 0.45, 12, 0.5, kicker, 14, RED, bold=True)
    add_text(s, 0.7, 0.85, 12.2, 0.9, title, 26, DARK, bold=True)
    return s


# ------------------------------------------------------------- Slide 1 capa
s = prs.slides.add_slide(BLANK)
band(s)
add_text(s, 0.9, 2.0, 11.5, 0.5, "Case iFood", 18, RED, bold=True)
add_text(s, 0.9, 2.7, 11.8, 1.8,
         "O problema não é quanto se gasta em cupom.\nÉ para quem se manda.", 38, DARK, bold=True)
add_text(s, 0.9, 4.8, 11.5, 1.2,
         "Medimos o quanto cada envio realmente gera de venda nova — e descobrimos que\n"
         "9 em cada 10 clientes recebem a oferta errada.", 17, GREY)
add_text(s, 0.9, 6.5, 11.5, 0.5,
         "17 mil clientes · 306 mil eventos · 76 mil envios em 6 campanhas · 30 dias", 12, GREY)

# --------------------------------------------------------- Slide 2 problema
s = content_slide("O problema", "Quase um terço do desconto sai do caixa sem gerar nada")
s.shapes.add_picture(str(FIG / "10_funnel.png"), Inches(0.6), Inches(1.9), width=Inches(7.2))
add_text(s, 8.1, 2.0, 4.8, 1.0, "30% dos cupons são usados\nsem o cliente ter visto\na oferta antes", 20, RED, bold=True)
add_text(s, 8.1, 3.6, 4.9, 3.2,
         "O cliente faz uma compra normal,\natinge o valor mínimo sem saber\nque havia uma oferta — e o sistema\n"
         "desconta automaticamente.\n\nÉ uma venda que aconteceria de\nqualquer forma, só que mais barata\npara o cliente "
         "e mais cara\npara a empresa.", 15, DARK)

# ------------------------------------------------------- Slide 3 quanto vale
s = content_slide("O que descobrimos", "Enviar compensa — mas o efeito depende muito da oferta")
s.shapes.add_picture(str(FIG / "05_ate.png"), Inches(0.6), Inches(2.0), width=Inches(6.8))
add_text(s, 7.7, 2.0, 5.3, 0.7, "Como medimos", 16, DARK, bold=True)
add_text(s, 7.7, 2.6, 5.3, 2.6,
         "Em cada campanha, 25% dos clientes\nnão receberam nada. Esse grupo mostra\n"
         "o que teria acontecido sem o envio.\n\nA diferença entre os dois grupos é a\nvenda que o envio causou.", 15, DARK)
add_text(s, 7.7, 4.8, 5.3, 0.9, "R$ 9,21 de venda nova\npor cliente, em 7 dias", 22, GREEN, bold=True)
add_text(s, 7.7, 6.0, 5.3, 0.8, "Já descontado o custo do cupom.\nDesconto rende 34% mais que BOGO.", 14, DARK)

# ---------------------------------------------------------- Slide 4 solução
s = content_slide("A solução", "Dá para saber, antes de enviar, em quem a oferta faz diferença")
s.shapes.add_picture(str(FIG / "08_policy_validation.png"), Inches(0.6), Inches(2.0), width=Inches(6.8))
add_text(s, 7.7, 2.0, 5.3, 0.7, "Testado em campanhas futuras", 16, DARK, bold=True)
add_text(s, 7.7, 2.6, 5.3, 2.4,
         "O modelo ordena os clientes por\nquanto a oferta deve render em cada um.\n\n"
         "No grupo apontado como prioridade,\no envio gerou R$ 21 por cliente.\nNo último grupo, zero.", 15, DARK)
add_text(s, 7.7, 5.0, 5.3, 1.0, "Hoje, só 10% recebem\na oferta certa para eles", 20, RED, bold=True)
add_text(s, 7.7, 6.2, 5.3, 0.7, "O maior ganho não é enviar menos —\né enviar melhor.", 14, DARK)

# ------------------------------------------- Slide 5 conclusão + impacto
s = content_slide("A conclusão", "Com o mesmo orçamento, 20% mais resultado")
s.shapes.add_picture(str(FIG / "09_business_impact.png"), Inches(0.5), Inches(1.9), width=Inches(6.4))

add_text(s, 7.2, 1.85, 5.8, 0.6, "O que aprendemos", 17, DARK, bold=True)
for i, (titulo, txt) in enumerate([
    ("1. Cupom funciona — mas não para todos.",
     "O efeito médio é positivo, porém vai de R$ 21\na zero dependendo da pessoa."),
    ("2. Medir “quem usou o cupom” engana.",
     "Quem gasta muito usa cupom de qualquer jeito.\nSó comparando com quem não recebeu dá\npara saber o que a oferta causou."),
    ("3. O dinheiro está em acertar o alvo.",
     "Não em cortar envios: 90% recebem a oferta\nerrada, e corrigir isso rende +20%."),
]):
    y = 2.45 + i * 1.35
    add_text(s, 7.2, y, 5.9, 0.4, titulo, 14, RED, bold=True)
    add_text(s, 7.2, y + 0.35, 5.9, 0.9, txt, 13, DARK)

add_text(s, 0.6, 6.35, 6.4, 1.0,
         "Por 1 milhão de envios: de R$ 10,3M para R$ 12,3M — ganho baseado numa\n"
         "diferença já medida (desconto rende mais que BOGO). O potencial maior depende\n"
         "de o modelo acertar a oferta de cada pessoa, e por isso precisa de teste A/B.", 12, GREY)
add_text(s, 7.2, 6.5, 5.9, 0.6, "Próximo passo: teste A/B com grupo de controle.", 14, DARK, bold=True)

prs.save(str(OUT))
print("wrote", OUT)
